"""Generic file conversion service.

Converts between common document and image formats using only pip-installable
libraries — no LibreOffice, no system dependencies, no external API. Coverage:

  Input      →  Supported output targets
  ─────────     ─────────────────────────
  xlsx       →  pdf, csv, txt
  csv        →  xlsx, pdf, txt
  docx       →  pdf, txt
  pdf        →  txt (uses PyMuPDF; preserves text order per page)
  pptx       →  txt, pdf
  txt / md   →  pdf, docx
  html       →  pdf  (lightweight — uses html2text + reportlab, not a browser)
  png/jpg/webp/gif  →  pdf, and any other image format in the set

Each public converter returns (output_bytes, output_mime, out_filename).
The dispatcher `convert_bytes()` accepts raw bytes + from_ext + to_ext and
picks the right pair.

This file is deliberately self-contained so it can be imported from the plan
executor (`run_plan`) or from a standalone HTTP endpoint.
"""

from __future__ import annotations

import asyncio
import csv
import io
import logging
import os
import shutil
import tempfile
import uuid
from typing import Tuple

logger = logging.getLogger(__name__)


# ── Helpers ─────────────────────────────────────────────────────────────────

def _norm_ext(ext: str) -> str:
    """Normalize a format hint like '.XLSX', 'xlsx', 'application/pdf' → 'xlsx'."""
    e = (ext or "").strip().lower().lstrip(".")
    alias = {
        "jpeg": "jpg",
        "text": "txt",
        "markdown": "md",
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "text/csv": "csv",
        "text/plain": "txt",
        "text/html": "html",
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/webp": "webp",
        "image/gif": "gif",
    }
    return alias.get(e, e)


def _mime_for(ext: str) -> str:
    """Pick a reasonable content type for an output extension."""
    return {
        "pdf": "application/pdf",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "csv": "text/csv",
        "txt": "text/plain",
        "md": "text/markdown",
        "html": "text/html",
        "png": "image/png",
        "jpg": "image/jpeg",
        "webp": "image/webp",
        "gif": "image/gif",
    }.get(ext, "application/octet-stream")


# ── PDF builder (shared by several targets) ─────────────────────────────────

def _pdf_from_rows(rows: list[list], title: str = "") -> bytes:
    """Render a 2D list as a paginated PDF table. Used by xlsx/csv → pdf."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=landscape(letter), leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()
    story = []
    if title:
        story.append(Paragraph(title, styles["Heading2"]))
        story.append(Spacer(1, 8))
    if not rows:
        story.append(Paragraph("(empty)", styles["BodyText"]))
    else:
        # Stringify everything for reportlab and cap very long cells
        data = [[(str(c) if c is not None else "")[:200] for c in r] for r in rows]
        table = Table(data, repeatRows=1)
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#111")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ]))
        story.append(table)
    doc.build(story)
    return buf.getvalue()


def _pdf_from_text(text: str, title: str = "") -> bytes:
    """Render plain text (with optional title) as a PDF."""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, leftMargin=54, rightMargin=54, topMargin=54, bottomMargin=54)
    styles = getSampleStyleSheet()
    story = []
    if title:
        story.append(Paragraph(title, styles["Heading2"]))
        story.append(Spacer(1, 10))
    for para in (text or "").split("\n\n"):
        # Escape reportlab's XML-ish Paragraph markup
        safe = (para.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))
        if safe.strip():
            story.append(Paragraph(safe.replace("\n", "<br/>"), styles["BodyText"]))
            story.append(Spacer(1, 6))
    doc.build(story)
    return buf.getvalue()


# ── xlsx ────────────────────────────────────────────────────────────────────

def _xlsx_to_rows(data: bytes) -> list[list]:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    ws = wb.active
    return [list(row) for row in ws.iter_rows(values_only=True)]


def xlsx_to_pdf(data: bytes) -> bytes:
    return _pdf_from_rows(_xlsx_to_rows(data), title="")


def xlsx_to_csv(data: bytes) -> bytes:
    buf = io.StringIO()
    writer = csv.writer(buf)
    for row in _xlsx_to_rows(data):
        writer.writerow(["" if c is None else c for c in row])
    return buf.getvalue().encode("utf-8")


def xlsx_to_txt(data: bytes) -> bytes:
    rows = _xlsx_to_rows(data)
    return "\n".join("\t".join("" if c is None else str(c) for c in r) for r in rows).encode("utf-8")


# ── csv ─────────────────────────────────────────────────────────────────────

def _csv_to_rows(data: bytes) -> list[list]:
    # Try utf-8 first, fall back to latin-1 for legacy exports
    try:
        text = data.decode("utf-8-sig")
    except UnicodeDecodeError:
        text = data.decode("latin-1", errors="replace")
    return [row for row in csv.reader(io.StringIO(text))]


def csv_to_xlsx(data: bytes) -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    for row in _csv_to_rows(data):
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def csv_to_pdf(data: bytes) -> bytes:
    return _pdf_from_rows(_csv_to_rows(data))


def csv_to_txt(data: bytes) -> bytes:
    rows = _csv_to_rows(data)
    return "\n".join("\t".join(r) for r in rows).encode("utf-8")


# ── docx ────────────────────────────────────────────────────────────────────

def _docx_to_text(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text]
    # Pull simple table text too so conversions don't silently drop tables
    for tbl in doc.tables:
        for row in tbl.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n\n".join(parts)


def docx_to_pdf(data: bytes) -> bytes:
    return _pdf_from_text(_docx_to_text(data))


def docx_to_txt(data: bytes) -> bytes:
    return _docx_to_text(data).encode("utf-8")


# ── pdf ─────────────────────────────────────────────────────────────────────

def pdf_to_txt(data: bytes) -> bytes:
    import fitz  # PyMuPDF
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        pages = [page.get_text("text") for page in doc]
    finally:
        doc.close()
    return "\n\n".join(pages).encode("utf-8")


# ── pptx ────────────────────────────────────────────────────────────────────

def _pptx_to_text(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        out.append(f"Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                out.append(shape.text)
        out.append("")
    return "\n".join(out)


def pptx_to_txt(data: bytes) -> bytes:
    return _pptx_to_text(data).encode("utf-8")


def pptx_to_pdf(data: bytes) -> bytes:
    return _pdf_from_text(_pptx_to_text(data), title="Presentation")


# ── txt / md ────────────────────────────────────────────────────────────────

def txt_to_pdf(data: bytes) -> bytes:
    return _pdf_from_text(data.decode("utf-8", errors="replace"))


def txt_to_docx(data: bytes) -> bytes:
    from docx import Document
    doc = Document()
    for para in data.decode("utf-8", errors="replace").split("\n\n"):
        doc.add_paragraph(para)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── html ────────────────────────────────────────────────────────────────────

def html_to_pdf(data: bytes) -> bytes:
    """Very simple HTML → PDF: strip tags, feed to reportlab. For pixel-perfect
    rendering we'd need weasyprint (system deps) or chromium — not today."""
    import re
    raw = data.decode("utf-8", errors="replace")
    # Collapse tags; preserve line breaks from <br>/<p>
    raw = re.sub(r"<br\s*/?>", "\n", raw, flags=re.I)
    raw = re.sub(r"</p\s*>", "\n\n", raw, flags=re.I)
    raw = re.sub(r"<[^>]+>", "", raw)
    raw = raw.replace("&nbsp;", " ").replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
    return _pdf_from_text(raw)


# ── images ──────────────────────────────────────────────────────────────────

_IMAGE_EXTS = {"png", "jpg", "webp", "gif"}


def image_to_image(data: bytes, to_ext: str) -> bytes:
    from PIL import Image
    img = Image.open(io.BytesIO(data))
    buf = io.BytesIO()
    pil_fmt = {"jpg": "JPEG", "png": "PNG", "webp": "WEBP", "gif": "GIF"}[to_ext]
    if pil_fmt == "JPEG" and img.mode in ("RGBA", "P"):
        img = img.convert("RGB")
    img.save(buf, format=pil_fmt)
    return buf.getvalue()


def image_to_pdf(data: bytes) -> bytes:
    from PIL import Image
    img = Image.open(io.BytesIO(data))
    if img.mode in ("RGBA", "P"):
        img = img.convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="PDF")
    return buf.getvalue()


# ── Dispatcher ──────────────────────────────────────────────────────────────

# Map (from_ext, to_ext) → converter function
_TABLE = {
    ("xlsx", "pdf"): xlsx_to_pdf,
    ("xlsx", "csv"): xlsx_to_csv,
    ("xlsx", "txt"): xlsx_to_txt,
    ("csv", "xlsx"): csv_to_xlsx,
    ("csv", "pdf"): csv_to_pdf,
    ("csv", "txt"): csv_to_txt,
    ("docx", "pdf"): docx_to_pdf,
    ("docx", "txt"): docx_to_txt,
    ("pdf", "txt"): pdf_to_txt,
    ("pptx", "txt"): pptx_to_txt,
    ("pptx", "pdf"): pptx_to_pdf,
    ("txt", "pdf"): txt_to_pdf,
    ("txt", "docx"): txt_to_docx,
    ("md", "pdf"): txt_to_pdf,
    ("md", "docx"): txt_to_docx,
    ("html", "pdf"): html_to_pdf,
}


def list_supported() -> list[str]:
    """Return a human-readable list of supported conversion pairs."""
    pairs = sorted(f"{a}→{b}" for (a, b) in _TABLE)
    # Image pairs
    for a in sorted(_IMAGE_EXTS):
        pairs.append(f"{a}→pdf")
        for b in sorted(_IMAGE_EXTS - {a}):
            pairs.append(f"{a}→{b}")
    return pairs


def convert_bytes(data: bytes, from_ext: str, to_ext: str, *, out_name: str = "output") -> Tuple[bytes, str, str]:
    """Convert `data` from `from_ext` to `to_ext`. Returns (bytes, mime, filename).

    Raises ValueError if the conversion pair isn't supported.

    This is the synchronous, pure-Python dispatcher. For LibreOffice-backed
    conversions (true pixel-preserving docx↔pdf, docx↔xlsx, pdf→docx, etc.)
    use `convert_bytes_async` instead.
    """
    src = _norm_ext(from_ext)
    dst = _norm_ext(to_ext)

    if src == dst:
        return data, _mime_for(dst), f"{out_name}.{dst}"

    # Image → image or image → pdf
    if src in _IMAGE_EXTS and dst in _IMAGE_EXTS:
        return image_to_image(data, dst), _mime_for(dst), f"{out_name}.{dst}"
    if src in _IMAGE_EXTS and dst == "pdf":
        return image_to_pdf(data), _mime_for("pdf"), f"{out_name}.pdf"

    fn = _TABLE.get((src, dst))
    if not fn:
        raise ValueError(f"Unsupported conversion: {src} → {dst}. Supported: {', '.join(list_supported())}")

    return fn(data), _mime_for(dst), f"{out_name}.{dst}"


# ── LibreOffice headless converter ─────────────────────────────────────────
#
# When the Docker image has `soffice` on PATH, we delegate high-fidelity
# office conversions to it. The native Python converters above stay as a
# fallback so local dev and older deployments keep working.

# Formats LibreOffice can read (from the common file type registry).
_LO_INPUT = {
    "doc", "docx", "odt", "rtf", "txt", "md", "html", "htm",
    "xls", "xlsx", "ods", "csv", "tsv",
    "ppt", "pptx", "odp",
    "pdf",  # import filter (read-only)
    "epub",
}
# Formats LibreOffice can write.
_LO_OUTPUT = {
    "pdf",
    "docx", "doc", "odt", "rtf", "txt", "html",
    "xlsx", "xls", "ods", "csv",
    "pptx", "ppt", "odp",
    "png", "jpg", "svg", "epub",
}

# Which (src, dst) pairs we PREFER to route through LibreOffice even if a
# native converter also exists. We use LO for anything where layout matters.
_LO_PREFERRED: set[tuple[str, str]] = {
    ("docx", "pdf"), ("doc", "pdf"), ("odt", "pdf"), ("rtf", "pdf"),
    ("pptx", "pdf"), ("ppt", "pdf"), ("odp", "pdf"),
    ("xlsx", "pdf"), ("xls", "pdf"), ("ods", "pdf"),
    ("html", "pdf"),
    # Office ↔ office (for editing flows)
    ("pdf", "docx"), ("pdf", "odt"), ("pdf", "txt"),
    ("docx", "odt"), ("odt", "docx"),
    ("xlsx", "ods"), ("ods", "xlsx"),
    ("pptx", "odp"), ("odp", "pptx"),
    # Cross-family edits (LO does a best-effort layout)
    ("docx", "html"), ("xlsx", "html"),
}

_LO_LOCK = asyncio.Lock()  # serialize soffice calls — the headless process
                           # is not safe to run fully concurrent with a
                           # shared user profile


def _soffice_available() -> bool:
    """Return True if soffice is on PATH. Cached via env lookup."""
    return shutil.which("soffice") is not None


def _lo_target_filter(dst: str) -> str:
    """Some target extensions need a specific LibreOffice export filter to
    force the modern OOXML output instead of the legacy binary format."""
    return {
        "docx": "docx:MS Word 2007 XML",
        "xlsx": "xlsx:Calc Office Open XML",
        "pptx": "pptx:Impress MS PowerPoint 2007 XML",
        "pdf": "pdf",
    }.get(dst, dst)


async def _convert_with_libreoffice(data: bytes, src_ext: str, dst_ext: str) -> bytes:
    """Run `soffice --headless --convert-to` in a unique temp dir.

    Each call gets its own user profile so multiple requests can safely
    share one process pool. The overall _LO_LOCK still serializes calls
    because soffice can still race on shared font caches if you hammer it.
    """
    if not _soffice_available():
        raise RuntimeError("LibreOffice (soffice) not available on this host")

    workdir = tempfile.mkdtemp(prefix="lo-")
    profile = os.path.join(workdir, "profile")
    outdir = os.path.join(workdir, "out")
    os.makedirs(profile, exist_ok=True)
    os.makedirs(outdir, exist_ok=True)
    in_name = f"input-{uuid.uuid4().hex}.{src_ext}"
    in_path = os.path.join(workdir, in_name)
    with open(in_path, "wb") as f:
        f.write(data)

    target = _lo_target_filter(dst_ext)
    cmd = [
        "soffice",
        "--headless",
        "--nologo",
        "--nofirststartwizard",
        "--norestore",
        f"-env:UserInstallation=file://{profile}",
        "--convert-to", target,
        "--outdir", outdir,
        in_path,
    ]

    async with _LO_LOCK:
        try:
            proc = await asyncio.create_subprocess_exec(
                *cmd,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            try:
                stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=120)
            except asyncio.TimeoutError:
                proc.kill()
                raise RuntimeError("LibreOffice conversion timed out after 120s")
            if proc.returncode != 0:
                raise RuntimeError(f"soffice exited {proc.returncode}: {stderr.decode('utf-8', 'replace')[:400]}")

            # soffice names the output after the input file with the new ext
            base_in = os.path.splitext(in_name)[0]
            out_path = os.path.join(outdir, f"{base_in}.{dst_ext}")
            if not os.path.exists(out_path):
                # Fall back to whatever file actually landed in outdir
                candidates = [f for f in os.listdir(outdir) if f.startswith(base_in)]
                if not candidates:
                    raise RuntimeError(f"soffice produced no output. stderr={stderr.decode('utf-8','replace')[:400]}")
                out_path = os.path.join(outdir, candidates[0])
            with open(out_path, "rb") as f:
                return f.read()
        finally:
            shutil.rmtree(workdir, ignore_errors=True)


async def convert_bytes_async(data: bytes, from_ext: str, to_ext: str, *, out_name: str = "output") -> Tuple[bytes, str, str]:
    """Async dispatcher. Prefers LibreOffice for office-format conversions
    when `soffice` is available; falls back to the pure-Python converters
    otherwise. Image conversions always stay on the Python path (Pillow is
    faster and more predictable than LO for bitmaps).
    """
    src = _norm_ext(from_ext)
    dst = _norm_ext(to_ext)

    if src == dst:
        return data, _mime_for(dst), f"{out_name}.{dst}"

    # Images — always Python/Pillow
    if src in _IMAGE_EXTS and (dst in _IMAGE_EXTS or dst == "pdf"):
        return convert_bytes(data, src, dst, out_name=out_name)

    use_lo = (
        _soffice_available()
        and src in _LO_INPUT
        and dst in _LO_OUTPUT
        and ((src, dst) in _LO_PREFERRED or (src, dst) not in _TABLE)
    )
    if use_lo:
        try:
            out = await _convert_with_libreoffice(data, src, dst)
            return out, _mime_for(dst), f"{out_name}.{dst}"
        except Exception as e:
            # Log and try the Python fallback if one exists
            logger.warning("LibreOffice conversion %s→%s failed: %s", src, dst, e)
            if (src, dst) not in _TABLE and not (src in _IMAGE_EXTS):
                raise

    # Fall through to the sync Python dispatcher
    return convert_bytes(data, src, dst, out_name=out_name)


def list_supported_lo() -> list[str]:
    """Return the LibreOffice-backed pairs we'll prefer when soffice exists."""
    return sorted(f"{a}→{b}" for (a, b) in _LO_PREFERRED)
